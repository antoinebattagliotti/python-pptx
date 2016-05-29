# encoding: utf-8

"""
Test suite for pptx.parts.slidelayout module
"""

from __future__ import absolute_import

import pytest

from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.parts.slidelayout import (
    _LayoutPlaceholders, _LayoutShapeFactory, _LayoutShapeTree,
    SlideLayoutPart
)
from pptx.parts.slidemaster import SlideMasterPart
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.slide import SlideLayout

from ..oxml.unitdata.shape import a_ph, a_pic, an_nvPr, an_nvSpPr, an_sp
from ..unitutil.cxml import element
from ..unitutil.mock import (
    class_mock, function_mock, instance_mock, method_mock
)


class DescribeSlideLayoutPart(object):

    def it_knows_the_slide_master_it_inherits_from(self, master_fixture):
        slide_layout, slide_master_ = master_fixture
        slide_master = slide_layout.slide_master
        slide_layout.part_related_by.assert_called_once_with(RT.SLIDE_MASTER)
        assert slide_master is slide_master_

    def it_provides_access_to_its_slide_layout(self, layout_fixture):
        slide_layout_part, SlideLayout_ = layout_fixture[:2]
        sldLayout, slide_layout_ = layout_fixture[2:]
        slide_layout = slide_layout_part.slide_layout
        SlideLayout_.assert_called_once_with(sldLayout, slide_layout_part)
        assert slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layout_fixture(self, SlideLayout_, slide_layout_):
        sldLayout = element('w:sldLayout')
        slide_layout_part = SlideLayoutPart(None, None, sldLayout)
        return slide_layout_part, SlideLayout_, sldLayout, slide_layout_

    @pytest.fixture
    def master_fixture(self, slide_master_, part_related_by_):
        slide_layout = SlideLayoutPart(None, None, None, None)
        return slide_layout, slide_master_

    # fixture components -----------------------------------

    @pytest.fixture
    def part_related_by_(self, request, slide_master_):
        return method_mock(
            request, SlideLayoutPart, 'part_related_by',
            return_value=slide_master_
        )

    @pytest.fixture
    def SlideLayout_(self, request, slide_layout_):
        return class_mock(
            request, 'pptx.parts.slidelayout.SlideLayout',
            return_value=slide_layout_
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMasterPart)


class Describe_LayoutShapeTree(object):

    def it_constructs_a_layout_placeholder_for_a_placeholder_shape(
            self, factory_fixture):
        layout_shapes, ph_elm_, _LayoutShapeFactory_, layout_placeholder_ = (
            factory_fixture
        )
        layout_placeholder = layout_shapes._shape_factory(ph_elm_)
        _LayoutShapeFactory_.assert_called_once_with(ph_elm_, layout_shapes)
        assert layout_placeholder is layout_placeholder_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def factory_fixture(
            self, ph_elm_, _LayoutShapeFactory_, layout_placeholder_):
        layout_shapes = _LayoutShapeTree(None)
        return (
            layout_shapes, ph_elm_, _LayoutShapeFactory_, layout_placeholder_
        )

    # fixture components ---------------------------------------------

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def _LayoutShapeFactory_(self, request, layout_placeholder_):
        return function_mock(
            request, 'pptx.parts.slidelayout._LayoutShapeFactory',
            return_value=layout_placeholder_
        )

    @pytest.fixture
    def ph_elm_(self, request):
        return instance_mock(request, CT_Shape)


class Describe_LayoutShapeFactory(object):

    def it_constructs_a_layout_placeholder_for_a_shape_element(
            self, factory_fixture):
        shape_elm, parent_, ShapeConstructor_, shape_ = factory_fixture
        shape = _LayoutShapeFactory(shape_elm, parent_)
        ShapeConstructor_.assert_called_once_with(shape_elm, parent_)
        assert shape is shape_

    # fixtures -------------------------------------------------------

    @pytest.fixture(params=['ph', 'sp', 'pic'])
    def factory_fixture(
            self, request, ph_bldr, slide_layout_, _LayoutPlaceholder_,
            layout_placeholder_, BaseShapeFactory_, base_shape_):
        shape_bldr, ShapeConstructor_, shape_mock = {
            'ph':  (ph_bldr, _LayoutPlaceholder_, layout_placeholder_),
            'sp':  (an_sp(), BaseShapeFactory_,   base_shape_),
            'pic': (a_pic(), BaseShapeFactory_,   base_shape_),
        }[request.param]
        shape_elm = shape_bldr.with_nsdecls().element
        return shape_elm, slide_layout_, ShapeConstructor_, shape_mock

    # fixture components -----------------------------------

    @pytest.fixture
    def BaseShapeFactory_(self, request, base_shape_):
        return function_mock(
            request, 'pptx.parts.slidelayout.BaseShapeFactory',
            return_value=base_shape_
        )

    @pytest.fixture
    def base_shape_(self, request):
        return instance_mock(request, BaseShape)

    @pytest.fixture
    def _LayoutPlaceholder_(self, request, layout_placeholder_):
        return class_mock(
            request, 'pptx.parts.slidelayout.LayoutPlaceholder',
            return_value=layout_placeholder_
        )

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def ph_bldr(self):
        return (
            an_sp().with_child(
                an_nvSpPr().with_child(
                    an_nvPr().with_child(
                        a_ph().with_idx(1))))
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayoutPart)


class Describe_LayoutPlaceholders(object):

    def it_constructs_a_layout_placeholder_for_a_placeholder_shape(
            self, factory_fixture):
        layout_placeholders, ph_elm_ = factory_fixture[:2]
        _LayoutShapeFactory_, layout_placeholder_ = factory_fixture[2:]
        layout_placeholder = layout_placeholders._shape_factory(ph_elm_)
        _LayoutShapeFactory_.assert_called_once_with(
            ph_elm_, layout_placeholders
        )
        assert layout_placeholder is layout_placeholder_

    def it_can_find_a_placeholder_by_idx_value(self, get_fixture):
        layout_placeholders, ph_idx, placeholder_ = get_fixture
        placeholder = layout_placeholders.get(idx=ph_idx)
        assert placeholder is placeholder_

    def it_returns_default_if_placeholder_having_idx_not_found(
            self, default_fixture):
        layout_placeholders = default_fixture
        default = 'barfoo'
        placeholder = layout_placeholders.get('foobar', default)
        assert placeholder is default

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def default_fixture(self, _iter_):
        layout_placeholders = _LayoutPlaceholders(None)
        return layout_placeholders

    @pytest.fixture
    def factory_fixture(
            self, ph_elm_, _LayoutShapeFactory_, layout_placeholder_):
        layout_placeholders = _LayoutPlaceholders(None)
        return (
            layout_placeholders, ph_elm_, _LayoutShapeFactory_,
            layout_placeholder_
        )

    @pytest.fixture(params=[0, 1])
    def get_fixture(self, request, _iter_, placeholder_, placeholder_2_):
        layout_placeholders = _LayoutPlaceholders(None)
        ph_idx = request.param
        ph_shape_ = {0: placeholder_, 1: placeholder_2_}[request.param]
        return layout_placeholders, ph_idx, ph_shape_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def _iter_(self, request, placeholder_, placeholder_2_):
        return method_mock(
            request, _LayoutPlaceholders, '__iter__',
            return_value=iter([placeholder_, placeholder_2_])
        )

    @pytest.fixture
    def layout_placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder)

    @pytest.fixture
    def _LayoutShapeFactory_(self, request, layout_placeholder_):
        return function_mock(
            request, 'pptx.parts.slidelayout._LayoutShapeFactory',
            return_value=layout_placeholder_
        )

    @pytest.fixture
    def ph_elm_(self, request):
        return instance_mock(request, CT_Shape)

    @pytest.fixture
    def placeholder_(self, request):
        return instance_mock(request, LayoutPlaceholder, idx=0)

    @pytest.fixture
    def placeholder_2_(self, request):
        return instance_mock(request, LayoutPlaceholder, idx=1)
